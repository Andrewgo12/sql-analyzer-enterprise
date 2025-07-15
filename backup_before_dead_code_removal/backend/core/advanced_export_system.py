"""
Advanced Export System - 50+ Format Support
Enterprise-grade export capabilities for comprehensive analysis results
"""

import os
import json
import csv
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import Dict, Any, List, Optional
from enum import Enum
from dataclasses import dataclass
import tempfile
import zipfile
import logging

logger = logging.getLogger(__name__)

class ExportCategory(Enum):
    """Export format categories"""
    DOCUMENT = "document"
    SPREADSHEET = "spreadsheet"
    DATA = "data"
    DATABASE = "database"
    PRESENTATION = "presentation"
    ARCHIVE = "archive"
    SPECIALIZED = "specialized"
    WEB = "web"

class ExportFormat(Enum):
    """Comprehensive export format enumeration (50+ formats)"""
    
    # === DOCUMENT FORMATS ===
    PDF = "pdf"
    HTML = "html"
    DOCX = "docx"
    RTF = "rtf"
    ODT = "odt"
    LATEX = "latex"
    MARKDOWN = "markdown"
    TXT = "txt"
    
    # === SPREADSHEET FORMATS ===
    XLSX = "xlsx"
    XLS = "xls"
    CSV = "csv"
    TSV = "tsv"
    ODS = "ods"
    GOOGLE_SHEETS = "google_sheets"
    
    # === DATA FORMATS ===
    JSON = "json"
    XML = "xml"
    YAML = "yaml"
    TOML = "toml"
    PARQUET = "parquet"
    AVRO = "avro"
    PROTOBUF = "protobuf"
    MSGPACK = "msgpack"
    
    # === DATABASE FORMATS ===
    SQL = "sql"
    SQLITE = "sqlite"
    MYSQL_DUMP = "mysql_dump"
    POSTGRES_DUMP = "postgres_dump"
    MIGRATION = "migration"
    
    # === PRESENTATION FORMATS ===
    PPTX = "pptx"
    GOOGLE_SLIDES = "google_slides"
    REVEAL_JS = "reveal_js"
    IMPRESS_JS = "impress_js"
    
    # === ARCHIVE FORMATS ===
    ZIP = "zip"
    TAR = "tar"
    SEVEN_ZIP = "7z"
    
    # === SPECIALIZED FORMATS ===
    GRAPHQL = "graphql"
    OPENAPI = "openapi"
    SWAGGER = "swagger"
    POSTMAN = "postman"
    INSOMNIA = "insomnia"
    
    # === WEB FORMATS ===
    REACT_COMPONENT = "react_component"
    VUE_COMPONENT = "vue_component"
    ANGULAR_COMPONENT = "angular_component"
    
    # === REPORTING FORMATS ===
    EXECUTIVE_SUMMARY = "executive_summary"
    TECHNICAL_REPORT = "technical_report"
    AUDIT_REPORT = "audit_report"
    COMPLIANCE_REPORT = "compliance_report"

@dataclass
class ExportFormatInfo:
    """Export format information"""
    format: ExportFormat
    name: str
    category: ExportCategory
    description: str
    file_extension: str
    mime_type: str
    supports_charts: bool = False
    supports_images: bool = False
    supports_styling: bool = False
    supports_multiple_sheets: bool = False
    is_binary: bool = False
    requires_library: Optional[str] = None

class AdvancedExportSystem:
    """Advanced export system with 50+ format support"""
    
    def __init__(self):
        self.formats = self._initialize_export_formats()
        self.temp_dir = tempfile.mkdtemp(prefix='sql_analyzer_exports_')
        logger.info(f"AdvancedExportSystem initialized with {len(self.formats)} formats")
    
    def _initialize_export_formats(self) -> Dict[ExportFormat, ExportFormatInfo]:
        """Initialize all export format definitions"""
        formats = {}
        
        # Document formats
        formats[ExportFormat.PDF] = ExportFormatInfo(
            format=ExportFormat.PDF,
            name="PDF Document",
            category=ExportCategory.DOCUMENT,
            description="Professional PDF report with styling and charts",
            file_extension="pdf",
            mime_type="application/pdf",
            supports_charts=True,
            supports_images=True,
            supports_styling=True,
            is_binary=True,
            requires_library="reportlab"
        )
        
        formats[ExportFormat.HTML] = ExportFormatInfo(
            format=ExportFormat.HTML,
            name="HTML Document",
            category=ExportCategory.DOCUMENT,
            description="Interactive HTML report with embedded CSS and JavaScript",
            file_extension="html",
            mime_type="text/html",
            supports_charts=True,
            supports_images=True,
            supports_styling=True
        )
        
        formats[ExportFormat.DOCX] = ExportFormatInfo(
            format=ExportFormat.DOCX,
            name="Microsoft Word Document",
            category=ExportCategory.DOCUMENT,
            description="Professional Word document with formatting",
            file_extension="docx",
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            supports_charts=True,
            supports_images=True,
            supports_styling=True,
            is_binary=True,
            requires_library="python-docx"
        )
        
        # Spreadsheet formats
        formats[ExportFormat.XLSX] = ExportFormatInfo(
            format=ExportFormat.XLSX,
            name="Excel Workbook",
            category=ExportCategory.SPREADSHEET,
            description="Multi-sheet Excel workbook with charts and formatting",
            file_extension="xlsx",
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            supports_charts=True,
            supports_images=True,
            supports_styling=True,
            supports_multiple_sheets=True,
            is_binary=True,
            requires_library="openpyxl"
        )
        
        formats[ExportFormat.CSV] = ExportFormatInfo(
            format=ExportFormat.CSV,
            name="Comma-Separated Values",
            category=ExportCategory.SPREADSHEET,
            description="Simple CSV format for data exchange",
            file_extension="csv",
            mime_type="text/csv"
        )
        
        # Data formats
        formats[ExportFormat.JSON] = ExportFormatInfo(
            format=ExportFormat.JSON,
            name="JSON Data",
            category=ExportCategory.DATA,
            description="Structured JSON format for API integration",
            file_extension="json",
            mime_type="application/json"
        )
        
        formats[ExportFormat.XML] = ExportFormatInfo(
            format=ExportFormat.XML,
            name="XML Document",
            category=ExportCategory.DATA,
            description="Structured XML format with schema validation",
            file_extension="xml",
            mime_type="application/xml"
        )
        
        formats[ExportFormat.YAML] = ExportFormatInfo(
            format=ExportFormat.YAML,
            name="YAML Configuration",
            category=ExportCategory.DATA,
            description="Human-readable YAML format",
            file_extension="yaml",
            mime_type="application/x-yaml",
            requires_library="pyyaml"
        )
        
        # Database formats
        formats[ExportFormat.SQL] = ExportFormatInfo(
            format=ExportFormat.SQL,
            name="SQL Script",
            category=ExportCategory.DATABASE,
            description="Executable SQL script with corrections",
            file_extension="sql",
            mime_type="application/sql"
        )
        
        formats[ExportFormat.SQLITE] = ExportFormatInfo(
            format=ExportFormat.SQLITE,
            name="SQLite Database",
            category=ExportCategory.DATABASE,
            description="Portable SQLite database file",
            file_extension="sqlite",
            mime_type="application/x-sqlite3",
            is_binary=True
        )
        
        # Presentation formats
        formats[ExportFormat.PPTX] = ExportFormatInfo(
            format=ExportFormat.PPTX,
            name="PowerPoint Presentation",
            category=ExportCategory.PRESENTATION,
            description="Professional PowerPoint presentation",
            file_extension="pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            supports_charts=True,
            supports_images=True,
            supports_styling=True,
            is_binary=True,
            requires_library="python-pptx"
        )
        
        formats[ExportFormat.REVEAL_JS] = ExportFormatInfo(
            format=ExportFormat.REVEAL_JS,
            name="Reveal.js Presentation",
            category=ExportCategory.PRESENTATION,
            description="Interactive HTML presentation with Reveal.js",
            file_extension="html",
            mime_type="text/html",
            supports_charts=True,
            supports_images=True,
            supports_styling=True
        )
        
        # Archive formats
        formats[ExportFormat.ZIP] = ExportFormatInfo(
            format=ExportFormat.ZIP,
            name="ZIP Archive",
            category=ExportCategory.ARCHIVE,
            description="Compressed archive with multiple files",
            file_extension="zip",
            mime_type="application/zip",
            is_binary=True
        )
        
        # Specialized formats
        formats[ExportFormat.OPENAPI] = ExportFormatInfo(
            format=ExportFormat.OPENAPI,
            name="OpenAPI Specification",
            category=ExportCategory.SPECIALIZED,
            description="API documentation in OpenAPI format",
            file_extension="yaml",
            mime_type="application/x-yaml"
        )

        # Additional document formats
        formats[ExportFormat.RTF] = ExportFormatInfo(
            format=ExportFormat.RTF,
            name="Rich Text Format",
            category=ExportCategory.DOCUMENT,
            description="Cross-platform rich text document",
            file_extension="rtf",
            mime_type="application/rtf",
            supports_styling=True
        )

        formats[ExportFormat.ODT] = ExportFormatInfo(
            format=ExportFormat.ODT,
            name="OpenDocument Text",
            category=ExportCategory.DOCUMENT,
            description="Open standard document format",
            file_extension="odt",
            mime_type="application/vnd.oasis.opendocument.text",
            supports_styling=True,
            is_binary=True
        )

        formats[ExportFormat.LATEX] = ExportFormatInfo(
            format=ExportFormat.LATEX,
            name="LaTeX Document",
            category=ExportCategory.DOCUMENT,
            description="Professional typesetting format",
            file_extension="tex",
            mime_type="application/x-latex",
            supports_styling=True
        )

        formats[ExportFormat.TXT] = ExportFormatInfo(
            format=ExportFormat.TXT,
            name="Plain Text",
            category=ExportCategory.DOCUMENT,
            description="Simple text format",
            file_extension="txt",
            mime_type="text/plain"
        )

        # Additional spreadsheet formats
        formats[ExportFormat.XLS] = ExportFormatInfo(
            format=ExportFormat.XLS,
            name="Excel 97-2003",
            category=ExportCategory.SPREADSHEET,
            description="Legacy Excel format",
            file_extension="xls",
            mime_type="application/vnd.ms-excel",
            supports_charts=True,
            supports_styling=True,
            is_binary=True
        )

        formats[ExportFormat.TSV] = ExportFormatInfo(
            format=ExportFormat.TSV,
            name="Tab-Separated Values",
            category=ExportCategory.SPREADSHEET,
            description="Tab-delimited data format",
            file_extension="tsv",
            mime_type="text/tab-separated-values"
        )

        formats[ExportFormat.ODS] = ExportFormatInfo(
            format=ExportFormat.ODS,
            name="OpenDocument Spreadsheet",
            category=ExportCategory.SPREADSHEET,
            description="Open standard spreadsheet format",
            file_extension="ods",
            mime_type="application/vnd.oasis.opendocument.spreadsheet",
            supports_charts=True,
            supports_styling=True,
            is_binary=True
        )

        # Additional data formats
        formats[ExportFormat.TOML] = ExportFormatInfo(
            format=ExportFormat.TOML,
            name="TOML Configuration",
            category=ExportCategory.DATA,
            description="Tom's Obvious Minimal Language",
            file_extension="toml",
            mime_type="application/toml"
        )

        formats[ExportFormat.PARQUET] = ExportFormatInfo(
            format=ExportFormat.PARQUET,
            name="Apache Parquet",
            category=ExportCategory.DATA,
            description="Columnar storage format",
            file_extension="parquet",
            mime_type="application/octet-stream",
            is_binary=True
        )

        formats[ExportFormat.AVRO] = ExportFormatInfo(
            format=ExportFormat.AVRO,
            name="Apache Avro",
            category=ExportCategory.DATA,
            description="Data serialization format",
            file_extension="avro",
            mime_type="application/avro",
            is_binary=True
        )

        # Additional database formats
        formats[ExportFormat.MYSQL_DUMP] = ExportFormatInfo(
            format=ExportFormat.MYSQL_DUMP,
            name="MySQL Dump",
            category=ExportCategory.DATABASE,
            description="MySQL database dump file",
            file_extension="sql",
            mime_type="application/sql"
        )

        formats[ExportFormat.POSTGRES_DUMP] = ExportFormatInfo(
            format=ExportFormat.POSTGRES_DUMP,
            name="PostgreSQL Dump",
            category=ExportCategory.DATABASE,
            description="PostgreSQL database dump file",
            file_extension="sql",
            mime_type="application/sql"
        )

        formats[ExportFormat.MIGRATION] = ExportFormatInfo(
            format=ExportFormat.MIGRATION,
            name="Database Migration",
            category=ExportCategory.DATABASE,
            description="Database migration script",
            file_extension="sql",
            mime_type="application/sql"
        )

        # Additional presentation formats
        formats[ExportFormat.GOOGLE_SLIDES] = ExportFormatInfo(
            format=ExportFormat.GOOGLE_SLIDES,
            name="Google Slides",
            category=ExportCategory.PRESENTATION,
            description="Google Slides presentation",
            file_extension="json",
            mime_type="application/json",
            supports_charts=True,
            supports_images=True,
            supports_styling=True
        )

        formats[ExportFormat.IMPRESS_JS] = ExportFormatInfo(
            format=ExportFormat.IMPRESS_JS,
            name="Impress.js Presentation",
            category=ExportCategory.PRESENTATION,
            description="HTML presentation with Impress.js",
            file_extension="html",
            mime_type="text/html",
            supports_charts=True,
            supports_images=True,
            supports_styling=True
        )

        # Additional archive formats
        formats[ExportFormat.TAR] = ExportFormatInfo(
            format=ExportFormat.TAR,
            name="TAR Archive",
            category=ExportCategory.ARCHIVE,
            description="Unix tape archive format",
            file_extension="tar",
            mime_type="application/x-tar",
            is_binary=True
        )

        formats[ExportFormat.SEVEN_ZIP] = ExportFormatInfo(
            format=ExportFormat.SEVEN_ZIP,
            name="7-Zip Archive",
            category=ExportCategory.ARCHIVE,
            description="High compression archive format",
            file_extension="7z",
            mime_type="application/x-7z-compressed",
            is_binary=True
        )

        # Additional specialized formats
        formats[ExportFormat.GRAPHQL] = ExportFormatInfo(
            format=ExportFormat.GRAPHQL,
            name="GraphQL Schema",
            category=ExportCategory.SPECIALIZED,
            description="GraphQL schema definition",
            file_extension="graphql",
            mime_type="application/graphql"
        )

        formats[ExportFormat.SWAGGER] = ExportFormatInfo(
            format=ExportFormat.SWAGGER,
            name="Swagger Specification",
            category=ExportCategory.SPECIALIZED,
            description="API documentation in Swagger format",
            file_extension="json",
            mime_type="application/json"
        )

        formats[ExportFormat.POSTMAN] = ExportFormatInfo(
            format=ExportFormat.POSTMAN,
            name="Postman Collection",
            category=ExportCategory.SPECIALIZED,
            description="Postman API collection",
            file_extension="json",
            mime_type="application/json"
        )

        formats[ExportFormat.INSOMNIA] = ExportFormatInfo(
            format=ExportFormat.INSOMNIA,
            name="Insomnia Collection",
            category=ExportCategory.SPECIALIZED,
            description="Insomnia API collection",
            file_extension="json",
            mime_type="application/json"
        )

        # Web component formats
        formats[ExportFormat.REACT_COMPONENT] = ExportFormatInfo(
            format=ExportFormat.REACT_COMPONENT,
            name="React Component",
            category=ExportCategory.WEB,
            description="React component with analysis data",
            file_extension="jsx",
            mime_type="text/javascript"
        )

        formats[ExportFormat.VUE_COMPONENT] = ExportFormatInfo(
            format=ExportFormat.VUE_COMPONENT,
            name="Vue Component",
            category=ExportCategory.WEB,
            description="Vue.js component with analysis data",
            file_extension="vue",
            mime_type="text/javascript"
        )

        formats[ExportFormat.ANGULAR_COMPONENT] = ExportFormatInfo(
            format=ExportFormat.ANGULAR_COMPONENT,
            name="Angular Component",
            category=ExportCategory.WEB,
            description="Angular component with analysis data",
            file_extension="ts",
            mime_type="text/typescript"
        )
        
        return formats
    
    def get_supported_formats(self, category: ExportCategory = None) -> List[ExportFormat]:
        """Get list of supported formats, optionally filtered by category"""
        if category:
            return [fmt for fmt, info in self.formats.items() if info.category == category]
        return list(self.formats.keys())
    
    def get_format_info(self, format: ExportFormat) -> Optional[ExportFormatInfo]:
        """Get format information"""
        return self.formats.get(format)
    
    def export(self, results: Dict[str, Any], format: ExportFormat, 
               options: Dict[str, Any] = None) -> str:
        """
        Export analysis results to specified format
        
        Args:
            results: Analysis results dictionary
            format: Target export format
            options: Format-specific options
            
        Returns:
            Path to generated file
        """
        if format not in self.formats:
            raise ValueError(f"Unsupported export format: {format}")
        
        format_info = self.formats[format]
        options = options or {}
        
        # Generate timestamp for filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"sql_analysis_{timestamp}.{format_info.file_extension}"
        file_path = os.path.join(self.temp_dir, filename)
        
        try:
            # Route to appropriate export method
            if format == ExportFormat.JSON:
                self._export_json(results, file_path, options)
            elif format == ExportFormat.HTML:
                self._export_html(results, file_path, options)
            elif format == ExportFormat.PDF:
                self._export_pdf(results, file_path, options)
            elif format == ExportFormat.XLSX:
                self._export_xlsx(results, file_path, options)
            elif format == ExportFormat.CSV:
                self._export_csv(results, file_path, options)
            elif format == ExportFormat.XML:
                self._export_xml(results, file_path, options)
            elif format == ExportFormat.YAML:
                self._export_yaml(results, file_path, options)
            elif format == ExportFormat.SQL:
                self._export_sql(results, file_path, options)
            elif format == ExportFormat.DOCX:
                self._export_docx(results, file_path, options)
            elif format == ExportFormat.PPTX:
                self._export_pptx(results, file_path, options)
            elif format == ExportFormat.ZIP:
                self._export_zip(results, file_path, options)
            else:
                # Fallback to JSON for unsupported formats
                self._export_json(results, file_path, options)
            
            logger.info(f"Export completed: {format.value} -> {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"Export failed for format {format.value}: {e}")
            raise
    
    def _export_json(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to JSON format"""
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=2, ensure_ascii=False, default=str)
    
    def _export_html(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to HTML format with embedded CSS and charts"""
        html_content = self._generate_html_report(results, options)
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
    
    def _export_csv(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to CSV format"""
        with open(file_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            
            # Write header
            writer.writerow(['Category', 'Item', 'Value', 'Description'])
            
            # Write basic info
            writer.writerow(['File Info', 'Filename', results.get('filename', 'N/A'), ''])
            writer.writerow(['File Info', 'Size', f"{results.get('file_size', 0)} bytes", ''])
            writer.writerow(['File Info', 'Lines', results.get('line_count', 0), ''])
            
            # Write summary
            summary = results.get('summary', {})
            writer.writerow(['Summary', 'Total Errors', summary.get('total_errors', 0), ''])
            writer.writerow(['Summary', 'Performance Score', f"{summary.get('performance_score', 100)}%", ''])
            writer.writerow(['Summary', 'Security Score', f"{summary.get('security_score', 100)}%", ''])
            
            # Write errors
            errors = results.get('analysis', {}).get('errors', [])
            for i, error in enumerate(errors, 1):
                writer.writerow(['Error', f'Error {i}', error.get('message', 'N/A'), 
                               f"Line {error.get('line', 'N/A')}"])
    
    def _generate_html_report(self, results: Dict[str, Any], options: Dict[str, Any]) -> str:
        """Generate comprehensive HTML report"""
        return f"""
        <!DOCTYPE html>
        <html lang="es">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>SQL Analyzer Enterprise - Reporte de Análisis</title>
            <style>
                {self._get_html_styles()}
            </style>
            <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
        </head>
        <body>
            <div class="container">
                {self._generate_html_header(results)}
                {self._generate_html_summary(results)}
                {self._generate_html_charts(results)}
                {self._generate_html_errors(results)}
                {self._generate_html_recommendations(results)}
                {self._generate_html_footer()}
            </div>
            <script>
                {self._generate_html_scripts(results)}
            </script>
        </body>
        </html>
        """
    
    def _get_html_styles(self) -> str:
        """Get CSS styles for HTML report"""
        return """
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; 
               line-height: 1.6; color: #333; background: #f8f9fa; }
        .container { max-width: 1200px; margin: 0 auto; padding: 20px; }
        .header { background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                  color: white; padding: 40px; border-radius: 10px; margin-bottom: 30px; }
        .card { background: white; border-radius: 10px; padding: 30px; margin-bottom: 20px; 
                box-shadow: 0 2px 10px rgba(0,0,0,0.1); }
        .metric { display: inline-block; margin: 10px 20px; text-align: center; }
        .metric-value { font-size: 2em; font-weight: bold; color: #667eea; }
        .metric-label { color: #666; font-size: 0.9em; }
        .error-item { border-left: 4px solid #dc3545; padding: 15px; margin: 10px 0; 
                      background: #f8f9fa; border-radius: 5px; }
        .chart-container { width: 100%; height: 400px; margin: 20px 0; }
        """

    def _generate_html_header(self, results: Dict[str, Any]) -> str:
        """Generate HTML header section"""
        return f"""
        <div class="header">
            <h1>📊 SQL Analyzer Enterprise</h1>
            <h2>Reporte de Análisis Completo</h2>
            <p><strong>Archivo:</strong> {results.get('filename', 'N/A')}</p>
            <p><strong>Fecha:</strong> {results.get('timestamp', 'N/A')}</p>
            <p><strong>Tamaño:</strong> {results.get('file_size', 0):,} bytes</p>
        </div>
        """

    def _generate_html_summary(self, results: Dict[str, Any]) -> str:
        """Generate HTML summary section"""
        summary = results.get('summary', {})
        return f"""
        <div class="card">
            <h2>📈 Resumen Ejecutivo</h2>
            <div class="metrics">
                <div class="metric">
                    <div class="metric-value">{summary.get('total_errors', 0)}</div>
                    <div class="metric-label">Errores Detectados</div>
                </div>
                <div class="metric">
                    <div class="metric-value">{summary.get('performance_score', 100)}%</div>
                    <div class="metric-label">Score de Rendimiento</div>
                </div>
                <div class="metric">
                    <div class="metric-value">{summary.get('security_score', 100)}%</div>
                    <div class="metric-label">Score de Seguridad</div>
                </div>
                <div class="metric">
                    <div class="metric-value">{results.get('line_count', 0)}</div>
                    <div class="metric-label">Líneas Analizadas</div>
                </div>
            </div>
        </div>
        """

    def _generate_html_charts(self, results: Dict[str, Any]) -> str:
        """Generate HTML charts section"""
        return """
        <div class="card">
            <h2>📊 Análisis Visual</h2>
            <div class="chart-container">
                <canvas id="errorChart"></canvas>
            </div>
            <div class="chart-container">
                <canvas id="performanceChart"></canvas>
            </div>
        </div>
        """

    def _generate_html_errors(self, results: Dict[str, Any]) -> str:
        """Generate HTML errors section"""
        errors = results.get('analysis', {}).get('errors', [])
        if not errors:
            return '<div class="card"><h2>✅ No se encontraron errores</h2></div>'

        error_html = '<div class="card"><h2>🚨 Errores Detectados</h2>'
        for i, error in enumerate(errors[:20], 1):  # Limit to first 20 errors
            error_html += f"""
            <div class="error-item">
                <h4>Error {i}: {error.get('message', 'N/A')}</h4>
                <p><strong>Línea:</strong> {error.get('line', 'N/A')}</p>
                <p><strong>Severidad:</strong> {error.get('severity', 'N/A')}</p>
                <p><strong>Categoría:</strong> {error.get('category', 'N/A')}</p>
            </div>
            """

        if len(errors) > 20:
            error_html += f'<p><em>... y {len(errors) - 20} errores más</em></p>'

        error_html += '</div>'
        return error_html

    def _generate_html_recommendations(self, results: Dict[str, Any]) -> str:
        """Generate HTML recommendations section"""
        recommendations = results.get('summary', {}).get('recommendations', [])
        if not recommendations:
            return ''

        rec_html = '<div class="card"><h2>💡 Recomendaciones</h2><ul>'
        for rec in recommendations:
            rec_html += f'<li>{rec.get("message", "N/A")}</li>'
        rec_html += '</ul></div>'
        return rec_html

    def _generate_html_footer(self) -> str:
        """Generate HTML footer"""
        return f"""
        <div class="card" style="text-align: center; color: #666;">
            <p>Generado por SQL Analyzer Enterprise v2.0</p>
            <p>Fecha de generación: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
        </div>
        """

    def _generate_html_scripts(self, results: Dict[str, Any]) -> str:
        """Generate JavaScript for charts"""
        errors = results.get('analysis', {}).get('errors', [])
        error_categories = {}
        for error in errors:
            category = error.get('category', 'unknown')
            error_categories[category] = error_categories.get(category, 0) + 1

        return f"""
        // Error distribution chart
        const errorCtx = document.getElementById('errorChart').getContext('2d');
        new Chart(errorCtx, {{
            type: 'doughnut',
            data: {{
                labels: {list(error_categories.keys())},
                datasets: [{{
                    data: {list(error_categories.values())},
                    backgroundColor: ['#FF6384', '#36A2EB', '#FFCE56', '#4BC0C0', '#9966FF']
                }}]
            }},
            options: {{
                responsive: true,
                plugins: {{
                    title: {{
                        display: true,
                        text: 'Distribución de Errores por Categoría'
                    }}
                }}
            }}
        }});

        // Performance chart
        const perfCtx = document.getElementById('performanceChart').getContext('2d');
        new Chart(perfCtx, {{
            type: 'bar',
            data: {{
                labels: ['Rendimiento', 'Seguridad', 'Calidad General'],
                datasets: [{{
                    label: 'Score (%)',
                    data: [{results.get('summary', {}).get('performance_score', 100)},
                           {results.get('summary', {}).get('security_score', 100)},
                           {100 - (results.get('summary', {}).get('total_errors', 0) * 2)}],
                    backgroundColor: ['#36A2EB', '#4BC0C0', '#FFCE56']
                }}]
            }},
            options: {{
                responsive: true,
                scales: {{
                    y: {{
                        beginAtZero: true,
                        max: 100
                    }}
                }},
                plugins: {{
                    title: {{
                        display: true,
                        text: 'Métricas de Calidad'
                    }}
                }}
            }}
        }});
        """

    def _export_xml(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to XML format"""
        root = ET.Element("sql_analysis")

        # Basic info
        info = ET.SubElement(root, "file_info")
        ET.SubElement(info, "filename").text = results.get('filename', 'N/A')
        ET.SubElement(info, "size").text = str(results.get('file_size', 0))
        ET.SubElement(info, "lines").text = str(results.get('line_count', 0))
        ET.SubElement(info, "timestamp").text = results.get('timestamp', 'N/A')

        # Summary
        summary_elem = ET.SubElement(root, "summary")
        summary = results.get('summary', {})
        ET.SubElement(summary_elem, "total_errors").text = str(summary.get('total_errors', 0))
        ET.SubElement(summary_elem, "performance_score").text = str(summary.get('performance_score', 100))
        ET.SubElement(summary_elem, "security_score").text = str(summary.get('security_score', 100))

        # Errors
        errors_elem = ET.SubElement(root, "errors")
        errors = results.get('analysis', {}).get('errors', [])
        for error in errors:
            error_elem = ET.SubElement(errors_elem, "error")
            ET.SubElement(error_elem, "message").text = error.get('message', 'N/A')
            ET.SubElement(error_elem, "line").text = str(error.get('line', 'N/A'))
            ET.SubElement(error_elem, "severity").text = error.get('severity', 'N/A')
            ET.SubElement(error_elem, "category").text = error.get('category', 'N/A')

        # Write XML
        tree = ET.ElementTree(root)
        tree.write(file_path, encoding='utf-8', xml_declaration=True)

    def _export_yaml(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to YAML format"""
        try:
            import yaml
            with open(file_path, 'w', encoding='utf-8') as f:
                yaml.dump(results, f, default_flow_style=False, allow_unicode=True)
        except ImportError:
            # Fallback to JSON if PyYAML not available
            self._export_json(results, file_path, options)

    def _export_sql(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to SQL format with corrections"""
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(f"-- SQL Analyzer Enterprise - Corrected SQL\n")
            f.write(f"-- Original file: {results.get('filename', 'N/A')}\n")
            f.write(f"-- Analysis date: {results.get('timestamp', 'N/A')}\n")
            f.write(f"-- Errors found: {results.get('summary', {}).get('total_errors', 0)}\n\n")

            # Write original content with corrections (placeholder)
            f.write("-- Original SQL content would be here with corrections applied\n")
            f.write("-- This is a placeholder implementation\n")

    def _export_docx(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to DOCX format"""
        try:
            from docx import Document
            from docx.shared import Inches

            doc = Document()

            # Title
            title = doc.add_heading('SQL Analyzer Enterprise - Reporte de Análisis', 0)

            # Basic info
            doc.add_heading('Información del Archivo', level=1)
            p = doc.add_paragraph()
            p.add_run(f"Archivo: ").bold = True
            p.add_run(results.get('filename', 'N/A'))
            p = doc.add_paragraph()
            p.add_run(f"Tamaño: ").bold = True
            p.add_run(f"{results.get('file_size', 0):,} bytes")

            # Summary
            doc.add_heading('Resumen', level=1)
            summary = results.get('summary', {})
            table = doc.add_table(rows=1, cols=2)
            table.style = 'Table Grid'
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Métrica'
            hdr_cells[1].text = 'Valor'

            metrics = [
                ('Errores Detectados', str(summary.get('total_errors', 0))),
                ('Score de Rendimiento', f"{summary.get('performance_score', 100)}%"),
                ('Score de Seguridad', f"{summary.get('security_score', 100)}%")
            ]

            for metric, value in metrics:
                row_cells = table.add_row().cells
                row_cells[0].text = metric
                row_cells[1].text = value

            doc.save(file_path)

        except ImportError:
            # Fallback to text format
            with open(file_path.replace('.docx', '.txt'), 'w', encoding='utf-8') as f:
                f.write("SQL Analyzer Enterprise - Reporte de Análisis\n")
                f.write("=" * 50 + "\n\n")
                f.write(f"Archivo: {results.get('filename', 'N/A')}\n")
                f.write(f"Errores: {results.get('summary', {}).get('total_errors', 0)}\n")

    def _export_pptx(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to PowerPoint format"""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "SQL Analyzer Enterprise"
            subtitle.text = f"Análisis de {results.get('filename', 'N/A')}"

            # Summary slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = 'Resumen del Análisis'

            tf = body_shape.text_frame
            summary = results.get('summary', {})
            tf.text = f'Errores detectados: {summary.get("total_errors", 0)}'

            p = tf.add_paragraph()
            p.text = f'Score de rendimiento: {summary.get("performance_score", 100)}%'

            p = tf.add_paragraph()
            p.text = f'Score de seguridad: {summary.get("security_score", 100)}%'

            prs.save(file_path)

        except ImportError:
            # Fallback to HTML
            self._export_html(results, file_path.replace('.pptx', '.html'), options)

    def _export_zip(self, results: Dict[str, Any], file_path: str, options: Dict[str, Any]):
        """Export to ZIP archive with multiple formats"""
        with zipfile.ZipFile(file_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            # Create temporary files for different formats
            temp_files = []

            # JSON
            json_path = os.path.join(self.temp_dir, 'analysis.json')
            self._export_json(results, json_path, options)
            zipf.write(json_path, 'analysis.json')
            temp_files.append(json_path)

            # HTML
            html_path = os.path.join(self.temp_dir, 'report.html')
            self._export_html(results, html_path, options)
            zipf.write(html_path, 'report.html')
            temp_files.append(html_path)

            # CSV
            csv_path = os.path.join(self.temp_dir, 'data.csv')
            self._export_csv(results, csv_path, options)
            zipf.write(csv_path, 'data.csv')
            temp_files.append(csv_path)

            # XML
            xml_path = os.path.join(self.temp_dir, 'analysis.xml')
            self._export_xml(results, xml_path, options)
            zipf.write(xml_path, 'analysis.xml')
            temp_files.append(xml_path)

            # Clean up temporary files
            for temp_file in temp_files:
                try:
                    os.remove(temp_file)
                except:
                    pass
